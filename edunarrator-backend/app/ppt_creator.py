from pptx import Presentation
from pptx.util import Inches
import os

def create_ppt_with_audio(slide_texts: list, audio_urls: list, output_path: str) -> str:
    prs = Presentation()
    slide_layout = prs.slide_layouts[1]  # Title + Content

    for i, (text, audio_url) in enumerate(zip(slide_texts, audio_urls)):
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = f"Slide {i+1}"
        slide.placeholders[1].text = text.strip()

        # NOTE: pptx does not support embedding audio directly.
        # Instead, add the audio URL as a note or placeholder for frontend/download use.
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = f"Audio URL: {audio_url}"

    prs.save(output_path)
    return output_path
